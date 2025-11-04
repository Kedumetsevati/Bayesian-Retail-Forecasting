# make_presentation.py
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Images you already created in this same folder
FORECAST_IMG = "bayes_forecast_plot2_forecast.png"
HISTORY_IMG  = "bayes_forecast_plot1_history.png"

OUT_PPTX = "Bayesian_Forecasting_for_Retail_Analysis.pptx"

# Color theme: Data Science Green
GREEN = RGBColor(0, 153, 102)
DARK  = RGBColor(40, 40, 40)

prs = Presentation()

def set_title_color(shape, color):
    try:
        shape.text_frame.paragraphs[0].runs[0].font.color.rgb = color
    except Exception:
        pass

def add_title_slide(title, subtitle, notes):
    slide = prs.slides.add_slide(prs.slide_layouts[0])  # Title
    slide.shapes.title.text = title
    set_title_color(slide.shapes.title, GREEN)
    slide.placeholders[1].text = subtitle
    slide.placeholders[1].text_frame.paragraphs[0].font.size = Pt(18)
    slide.placeholders[1].text_frame.paragraphs[0].font.color.rgb = DARK
    slide.notes_slide.notes_text_frame.text = notes

def add_text_slide(title, bullets, notes):
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content
    slide.shapes.title.text = title
    set_title_color(slide.shapes.title, GREEN)
    tf = slide.placeholders[1].text_frame
    tf.clear()
    for i, line in enumerate(bullets):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = line
        p.level = 0
        p.font.size = Pt(18)
        p.font.color.rgb = DARK
    slide.notes_slide.notes_text_frame.text = notes

def add_image_slide(title, img_path, caption, notes):
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title Only
    slide.shapes.title.text = title
    set_title_color(slide.shapes.title, GREEN)
    slide.shapes.add_picture(img_path, Inches(1.0), Inches(1.5), width=Inches(8.0))
    cap = slide.shapes.add_textbox(Inches(1.0), Inches(5.5), Inches(8.0), Inches(1.0))
    cap.text_frame.text = caption
    cap.text_frame.paragraphs[0].font.size = Pt(15)
    cap.text_frame.paragraphs[0].font.color.rgb = DARK
    slide.notes_slide.notes_text_frame.text = notes

# ---- Slides ----
add_title_slide(
    "Bayesian Forecasting for Retail Analysis",
    "Kedumetse Nadour Vati, PhD | drkedumvati@gmail.com | github.com/KedumetseVati | Edmonton, Alberta, Canada",
    "Introduce yourself, your PhD background, and that this project uses Bayesian forecasting to inform retail decisions."
)

add_text_slide(
    "Project Overview",
    [
        "Bayesian linear regression to forecast monthly retail revenue.",
        "Uncertainty modeled via posterior distributions (credible intervals).",
        "Goal: actionable, probabilistic forecasts for planning and budgeting."
    ],
    "Contrast probabilistic forecasts vs point estimates; emphasize business value."
)

add_text_slide(
    "Dataset & Features",
    [
        "36 months of synthetic Costco-style data (monthly Date).",
        "Features: Promo (Oct–Feb), Holiday (Nov–Dec), PriceIndex.",
        "Target: Revenue."
    ],
    "Note: synthetic but realistic; reproduces promotions and seasonality."
)

add_text_slide(
    "Methodology",
    [
        "Baseline: y = β0 + β1·t + β2·Promo + ε",
        "Extended: y = Baseline + β3·Holiday + β4·(PriceIndex−1) + ε",
        "Inference: MCMC (PyMC, NUTS) or conjugate posterior sampling.",
        "Evaluation: WAIC & LOO for model selection."
    ],
    "Briefly explain priors and why NUTS/MCMC; define credible intervals."
)

add_text_slide(
    "Model Comparison (WAIC & LOO)",
    [
        "Extended model outperforms baseline under WAIC & LOO.",
        "Holiday and PriceIndex materially improve predictive accuracy."
    ],
    "Frame results as evidence-based model selection."
)

add_image_slide(
    "Forecast Results",
    FORECAST_IMG,
    "6-month forecast (mean + 95% credible intervals). Seasonal peaks and promo effects are visible.",
    "Walk through interpreting the interval lines and uncertainty."
)

add_image_slide(
    "Historical Fit",
    HISTORY_IMG,
    "Posterior mean fit against 36 months of observed revenue.",
    "Comment on fit quality and any residual patterns."
)

add_text_slide(
    "Key Insights",
    [
        "Promotions lift sales by roughly +10–15%.",
        "Holiday months (Nov–Dec) show clear seasonal peaks.",
        "PriceIndex shows a negative elasticity with revenue.",
        "Forecast indicates continued growth with quantified uncertainty."
    ],
    "Highlight actionable business implications."
)

add_text_slide(
    "Streamlit App Summary",
    [
        "Interactive scenario testing (horizon, promo/holiday toggles, price trend).",
        "Instant updates to forecast table and chart."
    ],
    "This helps non-technical stakeholders run 'what-if' analyses."
)

add_text_slide(
    "Key Takeaways & Contact",
    [
        "Bayesian approach => probabilistic forecasts for better decisions.",
        "Feature-rich models improve accuracy vs. simple trend lines.",
        "End-to-end workflow: data → model → diagnostics → deployment.",
        "Contact: drkedumvati@gmail.com | github.com/KedumetseVati | Edmonton, Alberta, Canada"
    ],
    "Invite questions; emphasize rigor + practicality."
)

prs.save(OUT_PPTX)
print(f"Saved: {OUT_PPTX}")

